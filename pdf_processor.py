import PyPDF2
import io
from pptx import Presentation
import logging

def extract_text_from_file(uploaded_file):
    """
    Extract text from an uploaded file (PDF or PowerPoint).
    
    Args:
        uploaded_file: Streamlit UploadedFile object or Flask FileStorage object
    
    Returns:
        str: Extracted text from the file
    
    Raises:
        Exception: If file processing fails
    """
    logging.info("PDF DEBUG: Starting text extraction from file")
    
    try:
        # Handle both Streamlit and Flask file objects
        if hasattr(uploaded_file, 'filename') and uploaded_file.filename:
            # Flask FileStorage object
            filename = uploaded_file.filename
            logging.info(f"PDF DEBUG: Flask FileStorage object, filename: {filename}")
        elif hasattr(uploaded_file, 'name') and uploaded_file.name:
            # Streamlit UploadedFile object
            filename = uploaded_file.name
            logging.info(f"PDF DEBUG: Streamlit UploadedFile object, filename: {filename}")
        else:
            logging.error("PDF DEBUG: Invalid file object or no filename found")
            raise Exception("Invalid file object or no filename found")
        
        # Check if filename exists and has extension
        if not filename:
            logging.error("PDF DEBUG: No filename provided")
            raise Exception("No filename provided")
        
        # Get file extension
        if '.' not in filename:
            logging.error(f"PDF DEBUG: File has no extension. Filename: '{filename}'")
            raise Exception(f"File has no extension. Filename: '{filename}'")
        
        file_extension = filename.lower().split('.')[-1]
        logging.info(f"PDF DEBUG: File extension: {file_extension}")
        
        if file_extension == 'pdf':
            logging.info("PDF DEBUG: Processing PDF file")
            return extract_text_from_pdf(uploaded_file)
        elif file_extension == 'pptx':
            logging.info("PDF DEBUG: Processing PPTX file")
            return extract_text_from_pptx(uploaded_file)
        else:
            logging.error(f"PDF DEBUG: Unsupported file type: {file_extension}")
            raise Exception(f"Unsupported file type: {file_extension}")
            
    except Exception as e:
        logging.error(f"PDF DEBUG: Error in extract_text_from_file: {e}")
        logging.error(f"PDF DEBUG: Exception type: {type(e)}")
        import traceback
        logging.error(f"PDF DEBUG: Traceback: {traceback.format_exc()}")
        raise

def extract_text_from_pdf(uploaded_file):
    """
    Extract text from an uploaded PDF file using PyPDF2.
    
    Args:
        uploaded_file: Streamlit UploadedFile object or Flask FileStorage object
    
    Returns:
        str: Extracted text from the PDF
    
    Raises:
        Exception: If PDF processing fails
    """
    logging.info("PDF DEBUG: Starting PDF text extraction")
    
    try:
        # Handle both Streamlit and Flask file objects
        if hasattr(uploaded_file, 'getvalue'):
            # Streamlit UploadedFile object
            logging.info("PDF DEBUG: Reading from Streamlit UploadedFile object")
            pdf_bytes = io.BytesIO(uploaded_file.getvalue())
        else:
            # Flask FileStorage object
            logging.info("PDF DEBUG: Reading from Flask FileStorage object")
            pdf_bytes = io.BytesIO(uploaded_file.read())
        
        logging.info("PDF DEBUG: Creating PDF reader object")
        # Create a PDF reader object
        pdf_reader = PyPDF2.PdfReader(pdf_bytes)
        
        logging.info(f"PDF DEBUG: PDF has {len(pdf_reader.pages)} pages")
        
        # Initialize empty text string
        extracted_text = ""
        
        # Extract text from each page
        for page_num in range(len(pdf_reader.pages)):
            try:
                logging.info(f"PDF DEBUG: Processing page {page_num + 1}")
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()
                
                if page_text:
                    extracted_text += f"\n--- Page {page_num + 1} ---\n"
                    extracted_text += page_text
                    extracted_text += "\n"
                    logging.info(f"PDF DEBUG: Page {page_num + 1} processed successfully")
                else:
                    logging.warning(f"PDF DEBUG: Page {page_num + 1} had no text")
                    
            except Exception as e:
                logging.error(f"PDF DEBUG: Error processing page {page_num + 1}: {e}")
                continue  # Skip this page but continue with others
        
        logging.info(f"PDF DEBUG: Total extracted text length: {len(extracted_text)}")
        
        # Clean up the text
        logging.info("PDF DEBUG: Cleaning extracted text")
        extracted_text = clean_text(extracted_text)
        
        logging.info("PDF DEBUG: PDF text extraction completed successfully")
        return extracted_text
    
    except Exception as e:
        logging.error(f"PDF DEBUG: Critical error in PDF extraction: {e}")
        logging.error(f"PDF DEBUG: Exception type: {type(e)}")
        import traceback
        logging.error(f"PDF DEBUG: Traceback: {traceback.format_exc()}")
        raise Exception(f"Failed to extract text from PDF: {str(e)}")

def extract_text_from_pptx(uploaded_file):
    """
    Extract text from an uploaded PowerPoint file using python-pptx.
    
    Args:
        uploaded_file: Streamlit UploadedFile object or Flask FileStorage object
    
    Returns:
        str: Extracted text from the PowerPoint
    
    Raises:
        Exception: If PowerPoint processing fails
    """
    try:
        # Handle both Streamlit and Flask file objects
        if hasattr(uploaded_file, 'getvalue'):
            # Streamlit UploadedFile object
            pptx_bytes = io.BytesIO(uploaded_file.getvalue())
        else:
            # Flask FileStorage object
            pptx_bytes = io.BytesIO(uploaded_file.read())
        
        # Create a Presentation object
        prs = Presentation(pptx_bytes)
        
        # Initialize empty text string
        extracted_text = ""
        
        # Extract text from each slide
        for slide_num, slide in enumerate(prs.slides):
            extracted_text += f"\n--- Slide {slide_num + 1} ---\n"
            
            # Extract text from all text-containing shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    extracted_text += shape.text + "\n"
                    
                # Handle tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            row_text.append(cell.text.strip())
                        extracted_text += " | ".join(row_text) + "\n"
        
        # Clean up the text
        extracted_text = clean_text(extracted_text)
        
        return extracted_text
    
    except Exception as e:
        logging.error(f"PowerPoint extraction error: {e}")
        raise Exception(f"Failed to extract text from PowerPoint: {str(e)}")

def clean_text(text):
    """
    Clean and normalize extracted text.
    
    Args:
        text (str): Raw extracted text
    
    Returns:
        str: Cleaned text
    """
    if not text:
        return ""
    
    # Remove excessive whitespace
    lines = text.split('\n')
    cleaned_lines = []
    
    for line in lines:
        # Strip whitespace from each line
        line = line.strip()
        
        # Skip empty lines but keep page separators
        if line or line.startswith('---'):
            cleaned_lines.append(line)
    
    # Join lines back together
    cleaned_text = '\n'.join(cleaned_lines)
    
    # Remove excessive newlines (more than 2 consecutive)
    import re
    cleaned_text = re.sub(r'\n{3,}', '\n\n', cleaned_text)
    
    return cleaned_text

def get_pdf_metadata(uploaded_file):
    """
    Extract metadata from PDF file.
    
    Args:
        uploaded_file: Streamlit UploadedFile object or Flask FileStorage object
    
    Returns:
        dict: PDF metadata
    """
    try:
        if hasattr(uploaded_file, 'getvalue'):
            pdf_bytes = io.BytesIO(uploaded_file.getvalue())
        else:
            pdf_bytes = io.BytesIO(uploaded_file.read())
            
        pdf_reader = PyPDF2.PdfReader(pdf_bytes)
        
        metadata = {
            'num_pages': len(pdf_reader.pages),
            'title': pdf_reader.metadata.get('/Title', 'Unknown') if pdf_reader.metadata else 'Unknown',
            'author': pdf_reader.metadata.get('/Author', 'Unknown') if pdf_reader.metadata else 'Unknown',
            'subject': pdf_reader.metadata.get('/Subject', 'Unknown') if pdf_reader.metadata else 'Unknown',
            'creator': pdf_reader.metadata.get('/Creator', 'Unknown') if pdf_reader.metadata else 'Unknown',
        }
        
        return metadata
    
    except Exception as e:
        logging.error(f"Metadata extraction error: {e}")
        return {
            'num_pages': 0,
            'title': 'Unknown',
            'author': 'Unknown',
            'subject': 'Unknown',
            'creator': 'Unknown',
            'error': str(e)
        }
